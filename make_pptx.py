#!/usr/bin/env python
"""make_pptx.py -- build presentation.pptx from the project figures.

Theme: deep navy (#16233A) title/discussion slides, warm paper (#FBFAF7)
content slides, pathway accents blue/orange/aqua matching every figure.
Fonts: Georgia (titles) / Calibri (body) / Consolas (data) -- all stock on
Windows so nothing substitutes when it opens in PowerPoint."""
from pptx import Presentation
from pptx.util import Inches as IN, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

# single dark theme -- every slide sits on deep navy
NAVY   = RGBColor(0x16,0x23,0x3A)
NAVY2  = RGBColor(0x22,0x33,0x52)   # raised panel
LIGHT  = RGBColor(0xF5,0xF3,0xEE)   # primary text
LIGHT2 = RGBColor(0xB9,0xBD,0xC9)   # secondary text
SKY    = RGBColor(0x7F,0xB2,0xEE)   # blue accent readable on navy
ORANGE = RGBColor(0xF0,0x8A,0x5C)
AQUA   = RGBColor(0x3E,0xC9,0x96)
RED    = RGBColor(0xE6,0x67,0x67)
# aliases so the slide code below stays unchanged
PAPER, PANEL, INK, INK2, BLUE = NAVY, NAVY2, LIGHT, LIGHT2, SKY

W, H = IN(13.333), IN(7.5)
prs = Presentation(); prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    return tf

def runs(par, text, color, size, font="Calibri", bold=False, italic=False):
    """Mini-markup: **bold** and `mono`."""
    for tok in re.split(r"(\*\*.*?\*\*|`.*?`)", text):
        if not tok: continue
        r = par.add_run()
        if tok.startswith("**"):
            r.text = tok[2:-2]; r.font.bold = True
        elif tok.startswith("`"):
            r.text = tok[1:-1]; r.font.name = "Consolas"
            r.font.size = Pt(size-1); r.font.color.rgb = color
            r.font.bold = bold; continue
        else:
            r.text = tok; r.font.bold = bold
        r.font.name = font; r.font.size = Pt(size)
        r.font.color.rgb = color; r.font.italic = italic

def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.shadow.inherit = False
    return s

def header(slide, eyebrow, title, accent, on_navy=False):
    ink, sub = LIGHT, LIGHT2
    tf = box(slide, IN(0.7), IN(0.38), IN(12), IN(0.4))
    runs(tf.paragraphs[0], eyebrow.upper(), accent, 13, "Consolas", bold=True)
    tf2 = box(slide, IN(0.7), IN(0.72), IN(12), IN(0.9))
    runs(tf2.paragraphs[0], title, ink, 33, "Georgia", bold=True)
    rect(slide, IN(0.7), IN(1.52), IN(1.6), Pt(3.2), accent)

def footer(slide, n, total, on_navy=False):
    tf = box(slide, IN(10.6), IN(7.05), IN(2.4), IN(0.35))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    runs(p, f"Chokepoint neurons  ·  {n} / {total}", LIGHT2, 11, "Consolas")

def bullets(slide, items, x=IN(0.7), y=IN(1.85), w=IN(11.9), size=18,
            on_navy=False, gap=8):
    ink = LIGHT
    tf = box(slide, x, y, w, IN(7.3)-y)
    for i,(txt,opts) in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        c = opts.get("color", ink); s = opts.get("size", size)
        if opts.get("bullet", True):
            runs(p, "•  ", opts.get("accent", BLUE), s, bold=True)
        runs(p, txt, c, s, italic=opts.get("italic", False))
    return tf

def pic(slide, path, x, y, w=None, h=None):
    return slide.shapes.add_picture(path, x, y, width=w, height=h)

N = 12; k = [0]
def new(dark=False):
    s = prs.slides.add_slide(BLANK); bg(s, NAVY)
    k[0]+=1; footer(s, k[0], N, on_navy=dark)
    return s

# ---- 1 title (navy) --------------------------------------------------------
s = new(dark=True)
tf = box(s, IN(1.2), IN(2.0), IN(11), IN(0.4))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
runs(p, "FLYWIRE SUMMER INTERNSHIP 2026  ·  PROJECT 7", SKY, 14, "Consolas", bold=True)
tf = box(s, IN(1.2), IN(2.5), IN(11), IN(1.8))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
runs(p, "Chokepoint neurons across sensory pathways", LIGHT, 44, "Georgia", bold=True)
p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; p.space_before=Pt(6)
runs(p, "progress & framing check", LIGHT2, 24, "Georgia", italic=True)
tf = box(s, IN(1.2), IN(4.7), IN(11), IN(0.5))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
for label,col in [("● olfactory   ",BLUE),("● mechanosensory   ",ORANGE),("● visual",AQUA)]:
    runs(p, label, col, 16, bold=True)
tf = box(s, IN(1.2), IN(5.6), IN(11), IN(0.5))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
runs(p, "Pranav  ·  August 2026  ·  FAFB v783 + male CNS v0.9", LIGHT2, 14)

# ---- 2 proposal ------------------------------------------------------------
s = new(); header(s, "June · the proposal", "What I set out to find", BLUE)
# pipeline diagram
labels = [("senses\nsmell · touch · vision", LIGHT), ("interneurons\nthe middle of the brain", ORANGE),
          ("descending neurons\ncommands to the body", LIGHT)]
x = IN(1.5)
for i,(t,c) in enumerate(labels):
    sh = rect(s, x, IN(1.95), IN(2.9), IN(0.95), PANEL)
    sh.line.color.rgb = c; sh.line.width = Pt(1.75)
    tfd = sh.text_frame; tfd.word_wrap=True
    a,b = t.split("\n")
    runs(tfd.paragraphs[0], a, c, 16, bold=True); tfd.paragraphs[0].alignment=PP_ALIGN.CENTER
    p2=tfd.add_paragraph(); p2.alignment=PP_ALIGN.CENTER; runs(p2, b, INK2, 11)
    if i<2:
        tfa = box(s, x+IN(2.9), IN(2.1), IN(0.75), IN(0.6))
        pa=tfa.paragraphs[0]; pa.alignment=PP_ALIGN.CENTER; runs(pa,"→",INK2,28)
    x += IN(3.65)
bullets(s, [
 ("**Hypothesis:** some interneurons are bottlenecks — “if you removed them, information could no longer get through.”", {}),
 ("**Payoff question:** is any neuron a bottleneck for all three senses? That would be a genuinely critical hub of the whole brain.", {}),
 ("**Why it matters:** in human brain imaging, network hubs are the regions preferentially damaged in Alzheimer's and stroke. The fly is the only complete wiring diagram where the idea can be tested exactly.", {}),
 ("Both predictions get tested directly in this talk — and both turn out to be wrong in an interesting way.", {"color":INK2,"size":15,"italic":True,"bullet":False}),
], y=IN(3.2), gap=10)

# ---- 3 learning ------------------------------------------------------------
s = new(); header(s, "June – early July", "Learning the data", BLUE)
bullets(s, [
 ("Downloaded edge lists for five connectomes; learned the data by building neuron **embeddings** and exploring the graph structure.", {}),
 ("**Hop-size analysis** settled the pathway definition: 1 hop leaves too few real sensory→motor endpoints; 3 hops balloons past what exact path computation can handle. `2 hops` each way became the standard.", {}),
 ("First olfactory subgraph running end to end — sensory neurons to descending neurons as one directed graph.", {}),
], y=IN(2.0), gap=14)
stats = [("134,181","neurons, FAFB v783"),("2.7M","connected pairs"),("3","sensory pathways built")]
x = IN(0.9)
for num,lab in stats:
    sh = rect(s, x, IN(5.1), IN(3.4), IN(1.35), PANEL)
    tfd = sh.text_frame
    runs(tfd.paragraphs[0], num, SKY, 30, "Consolas", bold=True)
    p2=tfd.add_paragraph(); runs(p2, lab, INK2, 13)
    x += IN(3.9)

# ---- 4 DNp32 ---------------------------------------------------------------
s = new(); header(s, "mid-July · first result, first mistake", "DNp32 — a headline that wasn't real", RED)
bullets(s, [
 ("First full run used **unweighted** betweenness — every connection counted the same — and produced a headline: `DNp32`, a bottleneck shared by smell *and* vision.", {}),
 ("The problem: the edge list I had inherited carried **no synapse counts**. A 5-synapse whisper and a 500-synapse highway looked identical.", {}),
 ("With the real weighted Codex data, `DNp32`'s betweenness is **exactly zero** in both pathways. It sits on no strong route at all — a pure artifact.", {}),
], y=IN(1.95), gap=12)
sh = rect(s, IN(0.7), IN(4.85), IN(11.9), IN(1.7), PANEL)
rect(s, IN(0.7), IN(4.85), Pt(4.5), IN(1.7), RED)
tfd = box(s, IN(1.05), IN(5.0), IN(11.3), IN(1.5))
runs(tfd.paragraphs[0], "Weighting is not a refinement — it decides the answer.", INK, 17, bold=True)
p2=tfd.add_paragraph(); p2.space_before=Pt(4)
runs(p2, "An ablation isolating the weighting shows it alone changes `23–32 of each top 50`. After this, no result in the project went unreported without a control.", INK, 15)

# ---- 5 rebuilt -------------------------------------------------------------
s = new(); header(s, "early August", "Rebuilt on weighted data — the anatomy checks out", BLUE)
pic(s, "fig4_validated_olfactory.png",  IN(0.55), IN(1.85), w=IN(6.1))
pic(s, "fig4_validated_mechanosensory.png", IN(6.85), IN(1.85), w=IN(6.1))
bullets(s, [
 ("From connectivity alone, the ranking recovers textbook anatomy: olfactory bottlenecks concentrate in the **antennal lobe**, touch bottlenecks in the **gnathal ganglion**.", {}),
 ("Cross-modal check: **5 neurons** appear in two pathways' top 50, **none in all three** — the first crack in the hypothesis.", {"accent":ORANGE}),
], y=IN(5.35), size=16, gap=6)

# ---- 6 metric --------------------------------------------------------------
s = new(); header(s, "the week everything got tested", "Is “important” even well-defined?", ORANGE)
pic(s, "fig2_metric_disagreement.png", IN(2.4), IN(1.75), h=IN(3.35))
bullets(s, [
 ("Different importance metrics pick **almost completely different neurons**: betweenness vs. the BANC influence metric overlap at `1/50` and `0/50`.", {"accent":ORANGE}),
 ("Worse: in two of three pathways, betweenness largely reduces to **counting synapses** (23/50, 26/50 overlap with raw synapse count).", {"accent":ORANGE}),
 ("Rescue: a null that shuffles the network while **preserving every neuron's exact degree** — 200 trials, FDR-corrected — `37–45 of 50` survive. The top neurons matter because of *where they sit*, not how many synapses they have.", {"accent":ORANGE}),
], y=IN(5.25), size=15, gap=5)

# ---- 7 deletion ------------------------------------------------------------
s = new(); header(s, "the deletion test", "The hypothesis fails", RED)
pic(s, "fig3_deletion_damage.png", IN(2.55), IN(1.75), h=IN(3.15))
bullets(s, [
 ("Deleted each top neuron and re-solved every sensory→motor route: routes get **under 1% longer** on average, and **zero** source→target pairs are ever disconnected.", {"accent":RED}),
 ("The five cross-modal neurons are indistinguishable from ordinary top-50 neurons (z between −0.79 and +0.94).", {"accent":RED}),
 ("Growth note: I also caught my own scoring artifact here — deleting a neuron that is itself an endpoint trivially cuts its own pairs, which had inflated an earlier version of this result.", {"accent":RED}),
], y=IN(5.05), size=15, gap=5)

# ---- 8 replication ---------------------------------------------------------
s = new(); header(s, "replication · the strongest result", "It reproduces in a second brain", AQUA)
pic(s, "fig5_replication.png", IN(1.95), IN(1.7), h=IN(3.45))
bullets(s, [
 ("Entire pipeline rerun on the **Janelia male CNS** — different animal, different sex, different reconstruction pipeline.", {"accent":AQUA}),
 ("Cell-type rankings agree at `ρ = 0.45–0.71`; top-25 overlaps are 4.5–43× chance. The redundancy result and the degree null replicate too.", {"accent":AQUA}),
 ("Whatever betweenness measures here, it is a property of *the fly nervous system*, not of one reconstruction.", {"accent":AQUA}),
], y=IN(5.3), size=15, gap=5)

# ---- 9 ground truth --------------------------------------------------------
s = new(); header(s, "ground truth", "The physiology agrees — with the negative", SKY)
rows = [("","betweenness rank","experimental result"),
        ("T4, Mi9  (motion pathway)","~1000–2000","blocking → fly is motion-blind"),
        ("lLN2T_c  (my olfactory rank 1)","1","silencing → no measurable effect on odor coding")]
tbl = s.shapes.add_table(3, 3, IN(1.3), IN(1.9), IN(10.7), IN(1.9)).table
tbl.columns[0].width, tbl.columns[1].width, tbl.columns[2].width = IN(3.6), IN(2.5), IN(4.6)
for r,row in enumerate(rows):
    for c,val in enumerate(row):
        cell = tbl.cell(r,c); cell.fill.solid()
        cell.fill.fore_color.rgb = SKY if r==0 else (NAVY if r%2 else NAVY2)
        tfc = cell.text_frame; tfc.word_wrap=True
        runs(tfc.paragraphs[0], val, NAVY if r==0 else LIGHT,
             13 if r==0 else 15, "Consolas" if (c==1 or r==0) else "Calibri",
             bold=(r==0 or c==0))
bullets(s, [
 ("Pre-registered screen with matched controls: my validated types are **not** better characterized than mid-ranked neurons (2/24 vs 4/27) — and every necessity-grade result sits in the *control* group.", {"accent":SKY}),
 ("**Betweenness ranks traffic. Necessity is a property of populations.** T4 cells are individually redundant on any single path — exactly what scores low on betweenness and high on necessity when the whole population is blocked.", {"accent":SKY}),
], y=IN(4.35), size=16, gap=10)

# ---- 10 synthesis ----------------------------------------------------------
s = new(); header(s, "synthesis", "What I think the story is now", BLUE)
points = [("structured", BLUE, "The junctions are real: they beat degree-preserving nulls and reproduce across two animals."),
          ("but redundant", ORANGE, "No single point of failure anywhere — the pathway routes around every deletion."),
          ("and metric-dependent", AQUA, "Which neurons get called “critical” depends heavily on the metric you choose.")]
y = IN(2.1)
for word,col,rest in points:
    tf = box(s, IN(1.0), y, IN(11.3), IN(0.9))
    p = tf.paragraphs[0]
    runs(p, word+"   ", col, 24, "Georgia", bold=True)
    runs(p, rest, INK, 18)
    y += IN(1.15)
tf = box(s, IN(1.0), IN(5.7), IN(11.3), IN(1.0))
runs(tf.paragraphs[0], "The proposal's core prediction — remove a bottleneck and information stops — was tested directly and is wrong. I'd rather report that with controls than the reverse without them.", INK2, 15, italic=True)

# ---- 11 what's next --------------------------------------------------------
s = new(); header(s, "looking ahead", "What's to come", AQUA)
plans = [("Final presentation & poster", "Lead with the cross-dataset replication; frame around metric-dependence and structured redundancy. End of August."),
         ("Third dataset — BANC", "Weighted connections are already public; needs a sensory-class label mapping (~2 days). Would make every headline a three-connectome result."),
         ("Complete the null replication", "Male-CNS mechanosensory Null B (~6–8 h unattended) finishes the grid."),
         ("Release the pipeline", "Version-control the 25 scripts with a README — the analysis runs on any pathway subgraph, not just these three.")]
y = IN(1.95)
for i,(head,rest) in enumerate(plans,1):
    tf = box(s, IN(1.0), y, IN(11.3), IN(1.1))
    p = tf.paragraphs[0]
    runs(p, f"{i}.  ", AQUA, 20, "Consolas", bold=True)
    runs(p, head, LIGHT, 19, bold=True)
    p2 = tf.add_paragraph(); p2.space_before=Pt(2)
    runs(p2, "     "+rest, LIGHT2, 15)
    y += IN(1.22)

# ---- 12 backup -------------------------------------------------------------
s = new(); header(s, "backup", "In reserve for questions", INK2)
bullets(s, [
 ("**Null A** (weight permutation): real weights displace rankings further than permuted ones, z = +2.6 to +4.5, all pathways.", {"accent":INK2,"size":15}),
 ("**Cohort sweep:** deleting whole cell types is mostly additive (median 1.02–1.14× the sum of single deletions); a small minority (`VA2_adPN`, 3.9×) are genuinely more than their parts.", {"accent":INK2,"size":15}),
 ("**Current-flow betweenness abandoned:** 24 h unfinished on the *smallest* graph, and it requires an undirected graph — signal running backwards from motor to eye.", {"accent":INK2,"size":15}),
 ("**Motor endpoint (male CNS):** at 2 hops only mechanosensory reaches real motor neurons (4,961 sources); olfactory and visual reach essentially none — which retroactively justifies the descending-neuron endpoint.", {"accent":INK2,"size":15}),
 ("**Corrected-claims log:** DNp32 (unweighted artifact) · “zero pairs cut” (endpoint-scoring artifact) · 11× cohort claim (size vs. coherence confound) — each caught by the next control.", {"accent":INK2,"size":15}),
], y=IN(1.95), gap=10)

prs.save("presentation.pptx")
print("saved presentation.pptx,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
