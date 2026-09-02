#!/usr/bin/env python
"""make_talk_concise.py -- 5-slide, 5-minute problem/solution deck.

No step-by-step methods; one rigor line only. Usage:
    python3 make_talk_concise.py [navy|teal|light]   (default teal)
Outputs talk_concise[_theme].pptx.
"""
import os, re, sys
from pptx import Presentation
from pptx.util import Inches as IN, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.abspath(__file__))
THEME = (sys.argv[1] if len(sys.argv) > 1 else "teal").lower()
THEMES = {
    "navy": dict(INK=(0x0E,0x15,0x25), INK2=(0x1A,0x24,0x3B), INK3=(0x24,0x31,0x4E),
                 LIGHT=(0xF7,0xF5,0xF0), MUTE=(0x9A,0xA3,0xB5), SKY=(0x8A,0xB8,0xF0),
                 ORANGE=(0xF5,0x9E,0x6C), AQUA=(0x52,0xD2,0xA4), GOLD=(0xEF,0xC9,0x6B)),
    "teal": dict(INK=(0x0B,0x31,0x38), INK2=(0x11,0x40,0x49), INK3=(0x18,0x4F,0x5A),
                 LIGHT=(0xF6,0xF3,0xEC), MUTE=(0x9D,0xBC,0xC0), SKY=(0x8F,0xC9,0xF2),
                 ORANGE=(0xF6,0xA4,0x63), AQUA=(0x63,0xE0,0xB0), GOLD=(0xF2,0xCE,0x6E)),
    "light": dict(INK=(0xFA,0xF8,0xF3), INK2=(0xEF,0xEB,0xE1), INK3=(0xE6,0xE0,0xD2),
                  LIGHT=(0x21,0x2A,0x38), MUTE=(0x6E,0x77,0x85), SKY=(0x2E,0x63,0xA6),
                  ORANGE=(0xC2,0x5B,0x1F), AQUA=(0x0F,0x7A,0x5C), GOLD=(0x9A,0x6B,0x0F)),
}
T = THEMES[THEME]
def _c(k): return RGBColor(*T[k])
INK, INK2, INK3 = _c("INK"), _c("INK2"), _c("INK3")
LIGHT, MUTE = _c("LIGHT"), _c("MUTE")
SKY, ORANGE, AQUA, GOLD = _c("SKY"), _c("ORANGE"), _c("AQUA"), _c("GOLD")
SERIF, SANS = "Georgia", "Calibri"

W, H = IN(13.333), IN(7.5)
prs = Presentation(); prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]; MAIN = 5

def runs(par, text, color, size, font=SANS, bold=False, italic=False):
    for tok in re.split(r"(\*\*.*?\*\*)", text):
        if not tok: continue
        r = par.add_run()
        if tok.startswith("**"): r.text = tok[2:-2]; r.font.bold = True
        else: r.text = tok; r.font.bold = bold
        r.font.name = font; r.font.size = Pt(size)
        r.font.color.rgb = color; r.font.italic = italic

def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True; return tf

def rect(slide, x, y, w, h, color, line=None, lw=1.25):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.045; s.fill.solid(); s.fill.fore_color.rgb = color
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    s.shadow.inherit = False; return s

def bar(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False; return s

def new(accent=SKY):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = INK
    bar(s, 0, 0, W, IN(0.07), accent); return s

def header(slide, eyebrow, title, accent=SKY, tsize=28):
    bar(slide, IN(0.6), IN(0.52), IN(0.32), IN(0.055), accent)
    tf = box(slide, IN(1.05), IN(0.33), IN(11.6), IN(0.4))
    runs(tf.paragraphs[0], eyebrow.upper(), accent, 13, bold=True)
    tf = box(slide, IN(0.6), IN(0.72), IN(12.2), IN(1.0))
    runs(tf.paragraphs[0], title, LIGHT, tsize, font=SERIF, bold=True)

def footer(slide, n):
    tf = box(slide, IN(0.6), IN(7.06), IN(12.1), IN(0.35))
    runs(tf.paragraphs[0], "Chokepoint neurons across sensory pathways · FlyWire FAFB v783 + MCNS v0.9", MUTE, 10)
    tf2 = box(slide, IN(12.2), IN(7.06), IN(0.55), IN(0.35))
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    runs(p, f"{n} / {MAIN}", MUTE, 10, bold=True)

def bignum(slide, x, y, w, num, caption, accent, nsize=42, csize=13):
    tf = box(slide, x, y, w, IN(0.9))
    runs(tf.paragraphs[0], num, accent, nsize, font=SERIF, bold=True)
    tf = box(slide, x, y + IN(0.80), w, IN(0.75))
    runs(tf.paragraphs[0], caption, MUTE, csize)

def panel(slide, x, y, w, h, title, lines, accent, size=15):
    rect(slide, x, y, w, h, INK2)
    bar(slide, x, y + IN(0.14), IN(0.05), h - IN(0.28), accent)
    tf = box(slide, x + IN(0.22), y + IN(0.12), w - IN(0.4), IN(0.45))
    runs(tf.paragraphs[0], title, accent, 15, bold=True)
    tf = box(slide, x + IN(0.22), y + IN(0.55), w - IN(0.4), h - IN(0.62))
    for i, l in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5); runs(p, l, LIGHT, size)

def pic(slide, name, x, y, w=None, h=None):
    return slide.shapes.add_picture(os.path.join(ROOT, name), x, y, width=w, height=h)

# ============================================================ 1 · TITLE + PROBLEM HOOK
s = new(GOLD)
bar(s, 0, 0, IN(0.14), H, GOLD)
tf = box(s, IN(1.0), IN(0.85), IN(11.5), IN(0.5))
runs(tf.paragraphs[0], "FLYWIRE SUMMER INTERNSHIP 2026 · PROJECT 7", GOLD, 14, bold=True)
tf = box(s, IN(1.0), IN(1.4), IN(11.6), IN(2.2))
runs(tf.paragraphs[0], "Where does the fly's sensory traffic\nsqueeze through — and does anything\nbreak if you cut it?", LIGHT, 40, font=SERIF, bold=True)
tf = box(s, IN(1.0), IN(4.0), IN(11.5), IN(0.55))
runs(tf.paragraphs[0], "Finding — and stress-testing — the chokepoint neurons between the senses and the motor system", SKY, 19)
y0 = IN(4.8)
bignum(s, IN(1.0), y0, IN(2.6), "139k", "neurons · whole brain", SKY)
bignum(s, IN(3.7), y0, IN(2.6), "2.7M", "weighted connections", SKY)
bignum(s, IN(6.4), y0, IN(2.6), "3", "senses: smell · touch · vision", ORANGE)
bignum(s, IN(9.1), y0, IN(3.4), "2", "connectomes: FlyWire + male CNS", AQUA)
tf = box(s, IN(1.0), IN(6.5), IN(11.5), IN(0.5))
runs(tf.paragraphs[0], "Pranav Voleti · Princeton Neuroscience Institute symposium · August 2026", MUTE, 13)
footer(s, 1)

# ============================================================ 2 · THE PROBLEM
s = new(SKY)
header(s, "The problem", "We can now see every wire in the brain — but not which ones matter", SKY)
tf = box(s, IN(0.7), IN(1.95), IN(11.9), IN(1.7))
for i, l in enumerate([
    "Olfactory, mechanosensory and visual input must all converge onto the **descending neurons** — the shared endpoint every modality reaches within two hops.",
    "FAFB gives us the complete weighted graph. What it doesn't give us: **which neurons the sensory→DN traffic actually depends on.**",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12); runs(p, l, LIGHT, 19)
P, PH = IN(4.0), IN(2.2)
panel(s, IN(0.7),  P, IN(3.85), PH, "Which neurons?", ["rank highest by synapse-weighted betweenness on sensory→DN paths, per modality"], SKY, 16)
panel(s, IN(4.75), P, IN(3.85), PH, "Shared hubs?", ["are any high-betweenness neurons shared across modalities — brain-wide hubs?"], ORANGE, 16)
panel(s, IN(8.8),  P, IN(3.85), PH, "Truly critical?", ["does single-neuron deletion lengthen or sever sensory→DN paths — or does the graph route around it?"], AQUA, 16)
tf = box(s, IN(0.7), IN(6.45), IN(11.9), IN(0.5))
runs(tf.paragraphs[0], "Everything is replicated FAFB v783 → MaleCNS v0.9 at the cell-type level; nothing counts unless it recurs in both.", MUTE, 15, italic=True)
footer(s, 2)

# ============================================================ 3 · THE SOLUTION (main result)
s = new(ORANGE)
header(s, "The solution", "Each sense has its own chokepoints — and no single point of failure", ORANGE)
pic(s, "talk_fig_schematic.png", IN(1.27), IN(1.62), w=IN(10.8))
P2Y, P2H = IN(4.10), IN(1.78)
panel(s, IN(0.6), P2Y, IN(4.45), P2H, "1 · Where anatomy predicts", [
    "Olf → AL LNs (lLN2T_c, v2LN30, MZ_lv2PN)",
    "Mech → GNG DNs (DNge132, DNg62)",
    "Vis → LPi / Am1 / H2 wide-field cells",
], SKY, 15)
panel(s, IN(5.20), P2Y, IN(3.80), P2H, "2 · Conserved in a 2nd brain", [
    "**9 / 12 / 13** of top-25 types recur",
    "chance expectation: **< 1**",
], AQUA, 15)
panel(s, IN(9.15), P2Y, IN(3.60), P2H, "3 · None is required", [
    "delete any one → routes lengthen",
    "just **~0.5 %** · nothing disconnects",
], ORANGE, 15)
rect(s, IN(0.6), IN(6.02), IN(12.15), IN(0.88), INK3, line=ORANGE)
tf = box(s, IN(0.88), IN(6.12), IN(11.7), IN(0.8))
runs(tf.paragraphs[0], "Structured but redundant ", ORANGE, 18, bold=True)
runs(tf.paragraphs[0], " — conserved junctions carry the traffic, but every one has a parallel route.", LIGHT, 18)
footer(s, 3)

# ============================================================ 4 · PROOF IT HOLDS UP
s = new(AQUA)
header(s, "Does it hold up?", "Positioned beyond chance, replicated — but redundant", AQUA)
pic(s, "talk_fig_null.png", IN(0.5), IN(1.70), w=IN(8.0))
pic(s, "talk_fig_deletion.png", IN(8.75), IN(1.68), w=IN(4.15))
pic(s, "talk_fig_replication.png", IN(8.95), IN(4.08), w=IN(3.75))
rect(s, IN(0.6), IN(5.18), IN(7.75), IN(1.66), INK2)
bar(s, IN(0.6), IN(5.32), IN(0.05), IN(1.38), AQUA)
tf = box(s, IN(0.85), IN(5.30), IN(7.35), IN(1.5))
for i, l in enumerate([
    "Rankings tested against degree-preserving nulls (FDR q≤0.05) **plus a selection-bias control**: the full pipeline run on a shuffle-as-data graph, bounding the winner's-curse floor.",
    "Above the selection ceiling: **15 olfactory, 2 mechanosensory** neurons. Deletion damage spans **3 orders of magnitude** from top-50 (~0.5 %) to rank-1000 (~0.000 %).",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8); runs(p, l, LIGHT, 14.5)
footer(s, 4)

# ============================================================ 5 · TAKEAWAY
s = new(GOLD)
header(s, "Takeaway", "Conserved chokepoints, no single point of failure", GOLD)
y0 = IN(1.95)
bignum(s, IN(0.8), y0, IN(3.6), "9 / 12 / 13", "top-25 cell types conserved across two connectomes (chance < 1)", AQUA, 36)
bignum(s, IN(4.9), y0, IN(3.4), "~0.5 %", "route lengthening when any single chokepoint is deleted", ORANGE, 36)
bignum(s, IN(8.7), y0, IN(3.9), "6 / 200", "top neurons shared between two senses; none across all three", SKY, 36)
tf = box(s, IN(0.8), IN(3.9), IN(11.9), IN(1.7))
for i, l in enumerate([
    "**“Critical” depends on the question.**  Carrying the traffic, being necessary, and being conserved are three different things — and in this brain they disagree.",
    "**Next:** cohort-level deletions — T4/Mi9, behaviourally indispensable, sit at rank ~1000 in betweenness: necessity is a population property.",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(11); runs(p, "•  " + l, LIGHT if i == 0 else MUTE, 17)
rect(s, IN(0.6), IN(5.78), IN(12.15), IN(1.02), INK3, line=GOLD)
tf = box(s, IN(0.85), IN(5.90), IN(11.7), IN(0.9))
runs(tf.paragraphs[0], "The fly brain's sensory-motor chokepoints are real, they're conserved across two connectomes — and you can cut any one of them without breaking anything.", GOLD, 17, font=SERIF, italic=True)
footer(s, 5)

suffix = "" if THEME == "teal" else "_" + THEME
out = os.path.join(ROOT, f"talk_concise{suffix}.pptx")
prs.save(out); print("saved", out)
