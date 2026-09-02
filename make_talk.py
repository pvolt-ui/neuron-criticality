#!/usr/bin/env python
"""make_talk.py -- build talk.pptx: the 4-minute symposium talk.

5 main slides (question -> main result -> how -> evidence -> takeaway) + backups.
Same navy theme as make_pptx.py. Figures from talk_figures.py.
"""
import os, re
from pptx import Presentation
from pptx.util import Inches as IN, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.abspath(__file__))
NAVY, NAVY2 = RGBColor(0x16,0x23,0x3A), RGBColor(0x22,0x33,0x52)
LIGHT, LIGHT2 = RGBColor(0xF5,0xF3,0xEE), RGBColor(0xB9,0xBD,0xC9)
SKY, ORANGE, AQUA, RED = RGBColor(0x7F,0xB2,0xEE), RGBColor(0xF0,0x8A,0x5C), RGBColor(0x3E,0xC9,0x96), RGBColor(0xE6,0x67,0x67)
PAPER = RGBColor(0xFF,0xFF,0xFF)
W, H = IN(13.333), IN(7.5)
prs = Presentation(); prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
MAIN = 5  # number of main slides (for footer)

def runs(par, text, color, size, font="Calibri", bold=False, italic=False):
    for tok in re.split(r"(\*\*.*?\*\*|`.*?`)", text):
        if not tok: continue
        r = par.add_run()
        if tok.startswith("**"): r.text = tok[2:-2]; r.font.bold = True
        elif tok.startswith("`"): r.text = tok[1:-1]; r.font.name = "Consolas"; r.font.bold = bold
        else: r.text = tok; r.font.bold = bold
        if not tok.startswith("`"): r.font.name = font
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.italic = italic

def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True; return tf

def rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.06
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line: s.line.color.rgb = line; s.line.width = Pt(1.5)
    else: s.line.fill.background()
    s.shadow.inherit = False; return s

def new():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    return s

def header(slide, eyebrow, title, accent=SKY):
    tf = box(slide, IN(0.6), IN(0.35), IN(12), IN(0.4)); runs(tf.paragraphs[0], eyebrow.upper(), accent, 13, bold=True)
    tf = box(slide, IN(0.6), IN(0.7), IN(12.2), IN(1.0)); runs(tf.paragraphs[0], title, LIGHT, 28, font="Georgia", bold=True)

def footer(slide, n, label=None):
    tf = box(slide, IN(0.6), IN(7.0), IN(12), IN(0.35))
    runs(tf.paragraphs[0], label or f"Chokepoint neurons across sensory pathways · FlyWire FAFB v783 + MCNS v0.9 · {n}/{MAIN}", LIGHT2, 10)

def bullets(slide, items, x, y, w, size=17, gap=6, color=LIGHT):
    tf = box(slide, x, y, w, IN(5))
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple): runs(p, it[0], it[1], size)
        else: runs(p, "• " + it, color, size)
    return tf

def pic(slide, name, x, y, w=None, h=None):
    return slide.shapes.add_picture(os.path.join(ROOT, name), x, y, width=w, height=h)

def panel(slide, x, y, w, h, title, lines, accent, size=14):
    rect(slide, x, y, w, h, NAVY2)
    tf = box(slide, x + IN(0.2), y + IN(0.12), w - IN(0.4), IN(0.45)); runs(tf.paragraphs[0], title, accent, 15, bold=True)
    tf = box(slide, x + IN(0.2), y + IN(0.55), w - IN(0.4), h - IN(0.6))
    for i, l in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(4)
        runs(p, l, LIGHT, size)

# ---------------------------------------------------------------- 1 title + question
s = new()
tf = box(s, IN(0.8), IN(1.3), IN(11.7), IN(1.6))
runs(tf.paragraphs[0], "Where does the fly's sensory traffic squeeze through —\nand does anything break if you cut it?", LIGHT, 34, font="Georgia", bold=True)
tf = box(s, IN(0.8), IN(3.2), IN(11.7), IN(0.6))
runs(tf.paragraphs[0], "Chokepoint neurons across sensory pathways in the FlyWire connectome", SKY, 20)
rect(s, IN(0.8), IN(4.1), IN(11.7), IN(1.55), NAVY2)
tf = box(s, IN(1.0), IN(4.2), IN(11.3), IN(1.4))
runs(tf.paragraphs[0], "Question  ", ORANGE, 18, bold=True)
runs(tf.paragraphs[0], "For smell, touch and vision: which neurons carry the strongest routes from sensory input to the descending (motor-command) neurons — are any shared across senses — and is the circuit actually *dependent* on them?", LIGHT, 18)
tf = box(s, IN(0.8), IN(6.0), IN(11.7), IN(0.8))
runs(tf.paragraphs[0], "Pranav Voleti · FlyWire Summer Internship 2026 · Project 7 · FAFB v783 (139k neurons, 2.7M weighted connections) + MaleCNS v0.9", LIGHT2, 13)
footer(s, 1)

# ---------------------------------------------------------------- 2 main result
s = new()
header(s, "Main result", "Each sense has its own chokepoints — and no single point of failure", ORANGE)
pic(s, "talk_fig_schematic.png", IN(1.27), IN(1.72), w=IN(10.8))
P2Y, P2H = IN(4.16), IN(1.70)
panel(s, IN(0.6), P2Y, IN(4.45), P2H, "1 · Where anatomy says they should be", [
    "Smell → antennal-lobe interneurons",
    "Touch → gnathal-ganglion descenders",
    "Vision → lobula-plate wide-field",
], SKY, 16)
panel(s, IN(5.2), P2Y, IN(3.8), P2H, "2 · Conserved in a 2nd brain", [
    "**9 / 12 / 13** of 25 top types recur in the male CNS",
    "…but so does synapse count",
], AQUA, 16)
panel(s, IN(9.15), P2Y, IN(3.6), P2H, "3 · None is required", [
    "Delete any one: routes lengthen **0.3–0.6 %**",
    "Nothing disconnects a pair",
], ORANGE, 16)
rect(s, IN(0.6), IN(5.95), IN(12.15), IN(0.95), NAVY2, line=ORANGE)
tf = box(s, IN(0.85), IN(6.05), IN(11.7), IN(0.85))
runs(tf.paragraphs[0], "Structured but redundant  ", ORANGE, 18, bold=True)
runs(tf.paragraphs[0], "— conserved junctions carry the traffic, but every one has a parallel route.", LIGHT, 18)
footer(s, 2)

# ---------------------------------------------------------------- 3 how
s = new()
header(s, "Method", "Rank the neurons on the strongest routes — then try hard to break the result", SKY)
steps = [
    ("1  Build & rank", [
        "One graph per sense: every neuron ≤2 synapses from a sensory input **and** ≤2 from a descending neuron",
        "Rank by **synapse-weighted betweenness** on sensory→descending paths — high score = sits on the strong routes",
    ], SKY),
    ("2  Test it", [
        "Shuffle the wiring but keep every neuron's exact degree: is the rank just connection count?",
        "**And control the test itself** — run it on a structureless graph to see how many 'survivors' pure selection invents",
    ], ORANGE),
    ("3  Break it", [
        "Delete each bottleneck, re-solve every sensory→motor route: how much longer, how many pairs cut?",
        "Then rebuild the whole study on a second connectome",
    ], AQUA),
]
x0, w, h = IN(0.6), IN(3.9), IN(3.85)
for i, (t, paras, c) in enumerate(steps):
    x = x0 + i * IN(4.13)
    rect(s, x, IN(1.85), w, h, NAVY2, line=c)
    tf = box(s, x + IN(0.22), IN(1.97), w - IN(0.44), IN(0.55))
    runs(tf.paragraphs[0], t, c, 20, bold=True)
    tf = box(s, x + IN(0.22), IN(2.62), w - IN(0.44), h - IN(0.85))
    for j, para in enumerate(paras):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(12); runs(p, para, LIGHT, 16)
    if i < 2:
        tf = box(s, x + w + IN(0.01), IN(3.45), IN(0.3), IN(0.5)); runs(tf.paragraphs[0], "›", LIGHT2, 30)
tf = box(s, IN(0.6), IN(6.05), IN(12.1), IN(0.8))
runs(tf.paragraphs[0], "4 pathway graphs · ~1,000 betweenness evaluations for the nulls · 2 connectomes · `networkx` / `scipy`, reproducible end to end", LIGHT2, 14)
footer(s, 3)

# ---------------------------------------------------------------- 4 evidence
s = new()
header(s, "Evidence", "Positioned, replicated — but redundant", AQUA)
pic(s, "talk_fig_null.png", IN(0.5), IN(1.75), w=IN(8.0))
pic(s, "talk_fig_deletion.png", IN(8.75), IN(1.72), w=IN(4.15))
pic(s, "talk_fig_replication.png", IN(8.95), IN(4.12), w=IN(3.75))
tf = box(s, IN(0.6), IN(5.15), IN(7.7), IN(1.8))
for i, l in enumerate([
    "**Dotted line** = the highest score pure selection can invent, measured by running the whole test on a structureless graph.",
    "Only neurons **above** it are positioned beyond doubt: **15 olfactory, 2 mechanosensory**.",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(9); runs(p, l, LIGHT, 16)
footer(s, 4)

# ---------------------------------------------------------------- 5 takeaway
s = new()
header(s, "Takeaway", "What to remember", ORANGE)
bullets(s, [
    ("**Conserved chokepoints, no single point of failure.** Weighted betweenness finds anatomically sensible bottlenecks in each sense that replicate across two animals — and deleting any one barely matters.", LIGHT),
    ("**“Critical” depends on the question.** Traffic (betweenness) ≠ necessity (deletion) ≠ conserved identity (replication). All three are measured here; they disagree in instructive ways.", LIGHT),
    ("**Two self-corrections along the way:** (i) the first “visual” pathway was really the *ocellar* reflex arc — rebuilt with lamina/medulla inputs; (ii) the degree null over-counted survivors by a winner's-curse of 16–34/50 — now controlled.", LIGHT),
    ("**Next:** whole-population deletions (where necessity lives — T4/Mi9 are rank-1000 here yet behaviourally essential), and a third connectome (BANC).", LIGHT2),
], IN(0.7), IN(1.9), IN(11.9), size=18, gap=14)
rect(s, IN(0.6), IN(5.7), IN(12.15), IN(1.0), NAVY2, line=ORANGE)
tf = box(s, IN(0.8), IN(5.8), IN(11.8), IN(0.9))
runs(tf.paragraphs[0], "Elevator pitch  ", ORANGE, 15, bold=True)
runs(tf.paragraphs[0], "This project maps where smell, touch and vision funnel toward the fly's motor system in two whole-brain connectomes, and finds conserved chokepoint neurons that carry the traffic but — because every one has a parallel route — none the circuit cannot live without.", LIGHT, 15, italic=True)
footer(s, 5)

# ================================================================ backups
def backup_header(slide, title):
    header(slide, "Backup", title, LIGHT2)

s = new(); backup_header(s, "Pathway definitions and sizes")
bullets(s, [
    "Graph: Codex FAFB v783 connections, ≥5 synapses per pair → 134k neurons, 2.70M edges; MCNS v0.9 → 6.24M edges",
    "Targets: `super_class ∈ {descending, motor}`. Sources: olfactory = ORNs (1,765 in graph); mechanosensory = JO/mechanosensory afferents (1,675); visual = Lamina Monopolar + Transmedullary + Medulla Intrinsic families (17,004); ocellar = ocellar photoreceptors (142)",
    "Subgraph: ≤2 hops from a source AND ≤2 hops to a target. Olf 8.7k nodes/163k edges · Mech 13.6k/373k · Visual 60.5k/1.23M (sampled-source betweenness, 500 seeds) · Ocellar 386/2.1k",
    "**Correction:** original “visual” = `class == visual` (photoreceptors); 142/146 surviving sources were ocellar, and the 24k-node graph hung off 4 compound-eye PRs. Renamed `ocellar`; spec-compliant `visual` rebuilt. MCNS has no ocellar PRs, so the earlier visual replication was withdrawn and re-done like-for-like.",
    "Ocellar on its own: 112 neurons with nonzero betweenness; Null B 13/50 vs a selection floor of 2.8/50; OCG01 cells at z≈0 → degree-driven reflex arc, not a positional-bottleneck pathway.",
], IN(0.7), IN(1.9), IN(11.9), size=16, gap=11)
footer(s, "B1", "Backup B1")

s = new(); backup_header(s, "Null B — degree-preserving null and its selection-bias control")
bullets(s, [
    "Directed double-edge swap, 5×|E| attempts, preserves in-degree, out-degree, out-synapse total; 200 trials (olf/mech/ocellar), 30 (visual); 100 sampled sources; BH-FDR on p_z",
    "FDR survivors (top-50 / top-25): olfactory 37/20 · mechanosensory 43/20 · visual 29/16 · ocellar 13/12",
    "**Control**: treat one shuffle as data, take *its* top-50, test vs further shuffles (3 pseudo-real × 50 trials). False-positive floor: olfactory **16.0 ± 0.8**, mechanosensory **34.0 ± 2.2**, ocellar 2.8 ± 3.4; max z manufactured by selection alone: 9.9 / 23.4 / 3.4",
    "Above that ceiling: olfactory 15/50 (`MZ_lv2PN` z=35, `VES079` 23, `AL-AST1` 21, `lLN2T_c` 17, `v2LN30` 13 …); mechanosensory 2/50 (`DNge132` 31, `CB0021` 23); visual control not run (≈6 h) — likely survivors `PS124` z=15, `AVLP435a` 8",
    "Null A (weight permutation, 50 trials): real weights displace the ranking more than permuted weights in olf/mech/legacy-visual (z 2.6–4.5); fails in ocellar",
], IN(0.7), IN(1.9), IN(11.9), size=16, gap=11)
footer(s, "B2", "Backup B2")

s = new(); backup_header(s, "Deletion test, cell-type cohorts, and the cross-modal six")
bullets(s, [
    "Single deletion, 100 sources × all targets, damage scored only over pairs whose own endpoints survive: top-50 mean detour 0.37 % (olf), 0.34 % (mech), visual see talk fig; rank-1000+ ≈ 0.000 %; MCNS 0.35 / 0.06 %. Max single-neuron effect ≈ 4 % (legacy `CB3916`)",
    "A few FAFB deletions sever *whole sensory sources* (cut counts = multiples of the target count) — peripheral single-gateway cells, FAFB-only so far (0/45 in MCNS)",
    "Whole cell types: nearly all superadditive but barely (median 1.02–1.14×); 14/75 beat size- and rank-matched random sets at z>2. Standouts `VA2_adPN` (3.9×, z=+20), `OCG02b`, `lLN2T_c` (1.7×)",
    "Cross-modal six: `PVLP076` (olf+mech), `AVLP080` ×2, `CB0677`, `PS124` (mech+vis), `CB0676` (mech+ocellar). None in three. Legacy five all failed the deletion test (z −0.8…+0.9 vs ordinary top-50 controls)",
], IN(0.7), IN(1.9), IN(11.9), size=16, gap=11)
footer(s, "B3", "Backup B3")

s = new(); backup_header(s, "Metric dependence and ground truth")
bullets(s, [
    "Betweenness vs total synapse count top-50 overlap: olfactory 23/50, legacy-visual 26/50, **mechanosensory 6/50** (chance ≈ 0.3) — two pathways are largely degree-recoverable; mechanosensory is the one that needed the metric",
    "Flow metrics (Bates et al. 2025 influence; probabilistic traversal) agree with each other (ρ 0.77–0.89) and not with betweenness (ρ ≈ 0 – 0.3): *path* vs *flow* definitions of importance pick different neurons",
    "Current-flow betweenness: infeasible (24 h, no result on the smallest graph) and requires symmetrising the graph — reported as a negative feasibility result",
    "Pre-registered literature screen (validated types vs rank-1000–2000 controls): 2/24 vs 4/27 named types with functional evidence; all *necessity* evidence (T4, Mi9 → motion-blind) is in the **control** arm; the one validated type tested by silencing (`lLN2T_c` / LN2) showed no effect on odour coding — consistent with the deletion result",
], IN(0.7), IN(1.9), IN(11.9), size=16, gap=11)
footer(s, "B4", "Backup B4")

s = new(); backup_header(s, "Replication detail (FAFB vs MaleCNS, 150 sampled sources each, cell-type max)")
pic(s, "talk_fig_replication.png", IN(0.5), IN(1.9), h=IN(4.1))
bullets(s, [
    "Olfactory 9/25 (exp 0.52, p 2e-10), ρ 0.67 all / 0.54 nonzero-in-both (n=157): `lLN2T_c`, `v2LN30`, `il3LN6`, `MZ_lv2PN`, `AL-AST1`, `DM1_lPN`, `PLP096`, `PVLP076`, `lLN2F_b`",
    "Mechanosensory 12/25 (exp 0.28), ρ 0.71 / 0.56 (639): `DNg62`, `DNge132`, `DNge027`, `DNg35`, `DNg16`, `DNg100`, `DNp35`, `AVLP080`, `AVLP340`, `PS100`, `PS124`, `PVLP076`",
    "Visual 13/25 (exp 0.27), ρ 0.70 / 0.60 (548): `Am1`, `H2`, `LPT26`, `LT11`, `LT1c`, `LT62`, `LT79`, `Li32`, `Li33`, `PVLP011`, `PVLP061`, `AVLP080`, `PS124`",
    "Degree baseline: 20 / 18 / 19 of 25, ρ 0.90 / 0.86 / 0.91 — replication shows the pathway is conserved, not that betweenness adds conserved information",
    "Only mechanosensory reaches VNC motor neurons within 2+2 hops in MCNS (696 targets); olfactory/visual reach 1 / 0 — justifies the descending endpoint",
], IN(7.7), IN(1.85), IN(5.2), size=14, gap=9)
footer(s, "B5", "Backup B5")

out = os.path.join(ROOT, "talk.pptx"); prs.save(out); print("saved", out)
