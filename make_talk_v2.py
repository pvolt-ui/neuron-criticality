#!/usr/bin/env python
"""make_talk_v2.py -- build talk_v2.pptx: redesigned symposium deck.

5 main slides matching SCRIPT.md + 5 backups. Deep-ink theme, big-number
callouts, accent rules. Figures from talk_figures.py (already built).
"""
import os, re
from pptx import Presentation
from pptx.util import Inches as IN, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.abspath(__file__))

import sys
THEME = (sys.argv[1] if len(sys.argv) > 1 else "navy").lower()

THEMES = {
    # original deep-ink navy
    "navy": dict(
        INK=(0x0E,0x15,0x25), INK2=(0x1A,0x24,0x3B), INK3=(0x24,0x31,0x4E),
        LIGHT=(0xF7,0xF5,0xF0), MUTE=(0x9A,0xA3,0xB5),
        SKY=(0x8A,0xB8,0xF0), ORANGE=(0xF5,0x9E,0x6C),
        AQUA=(0x52,0xD2,0xA4), GOLD=(0xEF,0xC9,0x6B),
        OUT="talk_v2.pptx"),
    # dark, but deep teal-petrol instead of near-black
    "teal": dict(
        INK=(0x0B,0x31,0x38), INK2=(0x11,0x40,0x49), INK3=(0x18,0x4F,0x5A),
        LIGHT=(0xF6,0xF3,0xEC), MUTE=(0x9D,0xBC,0xC0),
        SKY=(0x8F,0xC9,0xF2),
        ORANGE=(0xF6,0xA4,0x63), AQUA=(0x63,0xE0,0xB0), GOLD=(0xF2,0xCE,0x6E),
        OUT="talk_teal.pptx"),
    # warm light "paper" theme
    "light": dict(
        INK=(0xFA,0xF8,0xF3), INK2=(0xEF,0xEB,0xE1), INK3=(0xE6,0xE0,0xD2),
        LIGHT=(0x21,0x2A,0x38), MUTE=(0x6E,0x77,0x85),
        SKY=(0x2E,0x63,0xA6), ORANGE=(0xC2,0x5B,0x1F),
        AQUA=(0x0F,0x7A,0x5C), GOLD=(0x9A,0x6B,0x0F),
        OUT="talk_light.pptx"),
}
T = THEMES[THEME]
def _c(t): return RGBColor(*t)
INK, INK2, INK3 = _c(T["INK"]), _c(T["INK2"]), _c(T["INK3"])
LIGHT, MUTE = _c(T["LIGHT"]), _c(T["MUTE"])
SKY, ORANGE, AQUA, GOLD = _c(T["SKY"]), _c(T["ORANGE"]), _c(T["AQUA"]), _c(T["GOLD"])

W, H = IN(13.333), IN(7.5)
prs = Presentation(); prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
MAIN = 5

SERIF, SANS = "Georgia", "Calibri"

def runs(par, text, color, size, font=SANS, bold=False, italic=False):
    for tok in re.split(r"(\*\*.*?\*\*|`.*?`)", text):
        if not tok: continue
        r = par.add_run()
        if tok.startswith("**"): r.text = tok[2:-2]; r.font.bold = True
        elif tok.startswith("`"): r.text = tok[1:-1]; r.font.name = "Consolas"
        else: r.text = tok; r.font.bold = bold
        if not tok.startswith("`"): r.font.name = font
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.italic = italic

def box(slide, x, y, w, h, anchor=None):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    if anchor: tf.vertical_anchor = anchor
    return tf

def rect(slide, x, y, w, h, color, line=None, radius=0.045, lw=1.25):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    s.shadow.inherit = False; return s

def bar(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False; return s

def new(accent=SKY):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = INK
    bar(s, 0, 0, W, IN(0.07), accent)          # top accent strip
    return s

def header(slide, eyebrow, title, accent=SKY, tsize=27):
    bar(slide, IN(0.6), IN(0.52), IN(0.32), IN(0.055), accent)
    tf = box(slide, IN(1.05), IN(0.33), IN(11.6), IN(0.4))
    runs(tf.paragraphs[0], eyebrow.upper(), accent, 13, bold=True)
    tf = box(slide, IN(0.6), IN(0.72), IN(12.2), IN(1.0))
    runs(tf.paragraphs[0], title, LIGHT, tsize, font=SERIF, bold=True)

def footer(slide, n, label=None):
    tf = box(slide, IN(0.6), IN(7.06), IN(12.1), IN(0.35))
    p = tf.paragraphs[0]
    runs(p, label or f"Chokepoint neurons across sensory pathways · FlyWire FAFB v783 + MCNS v0.9", MUTE, 10)
    tf2 = box(slide, IN(12.2), IN(7.06), IN(0.55), IN(0.35))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    runs(p2, f"{n} / {MAIN}" if isinstance(n, int) else str(n), MUTE, 10, bold=True)

def bignum(slide, x, y, w, num, caption, accent, nsize=44, csize=13, align=PP_ALIGN.LEFT):
    tf = box(slide, x, y, w, IN(0.9))
    p = tf.paragraphs[0]; p.alignment = align
    runs(p, num, accent, nsize, font=SERIF, bold=True)
    tf = box(slide, x, y + IN(0.82), w, IN(0.75))
    p = tf.paragraphs[0]; p.alignment = align
    runs(p, caption, MUTE, csize)

def panel(slide, x, y, w, h, title, lines, accent, size=15, tsize=15):
    rect(slide, x, y, w, h, INK2)
    bar(slide, x, y + IN(0.14), IN(0.05), h - IN(0.28), accent)
    tf = box(slide, x + IN(0.22), y + IN(0.12), w - IN(0.4), IN(0.45))
    runs(tf.paragraphs[0], title, accent, tsize, bold=True)
    tf = box(slide, x + IN(0.22), y + IN(0.55), w - IN(0.4), h - IN(0.62))
    for i, l in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        runs(p, l, LIGHT, size)

def bullets(slide, items, x, y, w, size=17, gap=10, color=LIGHT):
    tf = box(slide, x, y, w, IN(5))
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple): runs(p, it[0], it[1], size)
        else: runs(p, "•  " + it, color, size)
    return tf

def pic(slide, name, x, y, w=None, h=None):
    return slide.shapes.add_picture(os.path.join(ROOT, name), x, y, width=w, height=h)

# ================================================================ 1 · TITLE / QUESTION
s = new(GOLD)
bar(s, 0, 0, IN(0.14), H, GOLD)   # left spine
tf = box(s, IN(1.0), IN(0.9), IN(11.5), IN(0.5))
runs(tf.paragraphs[0], "FLYWIRE SUMMER INTERNSHIP 2026 · PROJECT 7", GOLD, 14, bold=True)
tf = box(s, IN(1.0), IN(1.45), IN(11.6), IN(2.2))
runs(tf.paragraphs[0], "Where does the fly's sensory traffic\nsqueeze through — and does anything\nbreak if you cut it?", LIGHT, 40, font=SERIF, bold=True)
tf = box(s, IN(1.0), IN(4.05), IN(11.5), IN(0.55))
runs(tf.paragraphs[0], "Chokepoint neurons on the sensory→motor pathways of two whole-brain connectomes", SKY, 19)

y0 = IN(4.85)
bignum(s, IN(1.0),  y0, IN(2.6), "139k",  "neurons · FAFB v783", SKY)
bignum(s, IN(3.7),  y0, IN(2.6), "2.7M",  "weighted connections", SKY)
bignum(s, IN(6.4),  y0, IN(2.6), "3",     "senses: smell · touch · vision", ORANGE)
bignum(s, IN(9.1),  y0, IN(3.4), "2",     "connectomes: FlyWire + male CNS", AQUA)

tf = box(s, IN(1.0), IN(6.55), IN(11.5), IN(0.5))
runs(tf.paragraphs[0], "Pranav Voleti · Princeton Neuroscience Institute symposium · August 2026", MUTE, 13)
footer(s, 1)

# ================================================================ 2 · MAIN RESULT
s = new(ORANGE)
header(s, "Main result", "Each sense has its own chokepoints — and no single point of failure", ORANGE)
pic(s, "talk_fig_schematic.png", IN(1.27), IN(1.62), w=IN(10.8))

P2Y, P2H = IN(4.10), IN(1.78)
panel(s, IN(0.6), P2Y, IN(4.45), P2H, "1 · Where anatomy predicts", [
    "Smell → antennal-lobe interneurons",
    "Touch → gnathal-ganglion descenders",
    "Vision → lobula-plate wide-field cells",
], SKY, 15)
panel(s, IN(5.20), P2Y, IN(3.80), P2H, "2 · Conserved in a 2nd brain", [
    "**9 / 12 / 13** of top-25 types recur",
    "chance expectation: **< 1**",
], AQUA, 15)
panel(s, IN(9.15), P2Y, IN(3.60), P2H, "3 · None is required", [
    "delete any one → routes lengthen",
    "just **~0.5 %** · nothing disconnects",
], ORANGE, 15)

rect(s, IN(0.6), IN(6.02), IN(12.15), IN(0.88), INK3, line=ORANGE)
tf = box(s, IN(0.88), IN(6.12), IN(11.7), IN(0.8))
runs(tf.paragraphs[0], "Structured but redundant ", ORANGE, 18, bold=True)
runs(tf.paragraphs[0], " — conserved junctions carry the traffic, but every one has a parallel route.", LIGHT, 18)
footer(s, 2)

# ================================================================ 3 · METHOD
s = new(SKY)
header(s, "Method", "Rank the neurons on the strongest routes — then try hard to break the result", SKY)
steps = [
    ("BUILD & RANK", SKY, [
        "One graph per sense: every neuron **≤2 synapses** from a sensory input *and* ≤2 from a descending neuron",
        "Rank by **synapse-weighted betweenness** — a neuron scores only if the *strong* routes run through it",
    ]),
    ("TEST IT — TWICE", ORANGE, [
        "Degree-preserving shuffle: is the ranking just connection count?",
        "**Control the test itself:** run everything on a structureless graph, where every 'hit' is a false positive by construction",
    ]),
    ("BREAK IT", AQUA, [
        "Delete each bottleneck; re-solve every sensory→motor route",
        "Rebuild the entire study in a second connectome",
    ]),
]
x0, w3, h3 = IN(0.6), IN(3.92), IN(3.75)
for i, (t, c, paras) in enumerate(steps):
    x = x0 + i * IN(4.12)
    rect(s, x, IN(1.85), w3, h3, INK2)
    bar(s, x, IN(1.85), w3, IN(0.09), c)
    tf = box(s, x + IN(0.24), IN(2.05), w3 - IN(0.48), IN(0.55))
    p = tf.paragraphs[0]
    runs(p, f"{i+1}", c, 30, font=SERIF, bold=True)
    runs(p, "   " + t, c, 17, bold=True)
    tf = box(s, x + IN(0.24), IN(2.75), w3 - IN(0.48), h3 - IN(1.0))
    for j, para in enumerate(paras):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(12); runs(p, para, LIGHT, 15.5)
    if i < 2:
        tf = box(s, x + w3 - IN(0.02), IN(3.35), IN(0.35), IN(0.6))
        runs(tf.paragraphs[0], "›", MUTE, 34, bold=True)
tf = box(s, IN(0.6), IN(5.95), IN(12.1), IN(0.8))
runs(tf.paragraphs[0], "4 pathway graphs · ~1,000 betweenness evaluations for the nulls · 2 connectomes · fully reproducible (`networkx` / `scipy`)", MUTE, 14)
footer(s, 3)

# ================================================================ 4 · EVIDENCE
s = new(AQUA)
header(s, "Evidence", "Positioned, replicated — but redundant", AQUA)
pic(s, "talk_fig_null.png", IN(0.5), IN(1.70), w=IN(8.0))
pic(s, "talk_fig_deletion.png", IN(8.75), IN(1.68), w=IN(4.15))
pic(s, "talk_fig_replication.png", IN(8.95), IN(4.08), w=IN(3.75))
rect(s, IN(0.6), IN(5.18), IN(7.75), IN(1.66), INK2)
bar(s, IN(0.6), IN(5.32), IN(0.05), IN(1.38), AQUA)
tf = box(s, IN(0.85), IN(5.30), IN(7.35), IN(1.5))
for i, l in enumerate([
    "**Dotted line** = the best score pure selection can invent (whole test run on a structureless graph).",
    "Only neurons **above** it are positioned beyond doubt: **15 olfactory · 2 mechanosensory**.",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8); runs(p, l, LIGHT, 15.5)
footer(s, 4)

# ================================================================ 5 · TAKEAWAY
s = new(GOLD)
header(s, "Takeaway", "Conserved chokepoints, no single point of failure", GOLD)
y0 = IN(1.95)
bignum(s, IN(0.8),  y0, IN(3.6), "9 / 12 / 13", "top-25 cell types conserved across two connectomes (chance < 1)", AQUA, 36)
bignum(s, IN(4.9),  y0, IN(3.4), "~0.5 %", "route lengthening when any single chokepoint is deleted", ORANGE, 36)
bignum(s, IN(8.7),  y0, IN(3.9), "6 / 200", "top neurons shared between two senses; none across all three", SKY, 36)

bullets(s, [
    ("**“Critical” depends on the question.**  Traffic (betweenness) ≠ necessity (deletion) ≠ conserved identity (replication) — measured here, and they disagree.", LIGHT),
    ("**Two self-corrections:** the first “visual” pathway was really the ocellar reflex arc; a selection-bias control caught inflation in the degree null.", LIGHT),
    ("**Next:** population-level deletions — behaviourally essential cells (T4, Mi9) sit at rank ~1000 here, so necessity lives at the population level.", MUTE),
], IN(0.8), IN(3.85), IN(11.9), size=16.5, gap=11)

rect(s, IN(0.6), IN(5.78), IN(12.15), IN(1.02), INK3, line=GOLD)
tf = box(s, IN(0.85), IN(5.90), IN(11.7), IN(0.9))
runs(tf.paragraphs[0], "The fly brain's sensory-motor chokepoints are real, they're conserved across two connectomes — and you can cut any one of them without breaking anything.", GOLD, 17, font=SERIF, italic=True)
footer(s, 5)

# ================================================================ BACKUPS (content unchanged from v1)
def backup(title, items, fig=None):
    s = new(MUTE); header(s, "Backup", title, MUTE, tsize=24)
    if fig:
        pic(s, fig, IN(0.5), IN(1.9), h=IN(4.1))
        bullets(s, items, IN(7.7), IN(1.85), IN(5.2), size=13.5, gap=8)
    else:
        bullets(s, items, IN(0.7), IN(1.9), IN(11.9), size=15.5, gap=10)
    return s

backup("Pathway definitions and sizes", [
    "Graph: Codex FAFB v783 connections, ≥5 synapses per pair → 134k neurons, 2.70M edges; MCNS v0.9 → 6.24M edges",
    "Targets: `super_class ∈ {descending, motor}`. Sources: olfactory = ORNs (1,765); mechanosensory = JO/mech afferents (1,675); visual = Lamina Monopolar + Tm + Mi families (17,004); ocellar = ocellar PRs (142)",
    "Subgraph ≤2 hops each side. Olf 8.7k/163k · Mech 13.6k/373k · Visual 60.5k/1.23M (sampled-source, 500 seeds) · Ocellar 386/2.1k",
    "**Correction:** original “visual” was 142 ocellar + 4 compound-eye PRs → renamed `ocellar`; spec-compliant `visual` rebuilt; MCNS visual replication withdrawn and redone like-for-like",
    "Ocellar alone: Null B 13/50 vs selection floor 2.8/50; OCG01 at z≈0 → degree-driven reflex arc, not positional",
]); footer(prs.slides[-1], "B1", "Backup B1")

backup("Null B — degree-preserving null and its selection-bias control", [
    "Directed double-edge swap, 5×|E| attempts, preserves in/out-degree and out-synapse total; 200 trials (olf/mech/ocellar), 30 (visual); BH-FDR on p_z",
    "FDR survivors (top-50 / top-25): olfactory 37/20 · mechanosensory 43/20 · visual 29/16 · ocellar 13/12",
    "**Control:** treat one shuffle as data, run the whole pipeline on it. False-positive floor: olf **16.0 ± 0.8**, mech **34.0 ± 2.2**, ocellar 2.8; max z from selection alone: 9.9 / 23.4 / 3.4",
    "Above the ceiling: olfactory 15/50 (`MZ_lv2PN` z=35, `VES079` 23, `AL-AST1` 21, `lLN2T_c` 17); mechanosensory 2/50 (`DNge132` 31, `CB0021` 23)",
    "Null A (weight permutation): real weights displace rankings more than permuted (z 2.6–4.5); fails in ocellar",
]); footer(prs.slides[-1], "B2", "Backup B2")

backup("Deletion test, cell-type cohorts, and the cross-modal six", [
    "Single deletion, 100 sources × all targets: top-50 mean detour 0.37 % (olf), 0.34 % (mech); rank-1000+ ≈ 0.000 %; MCNS 0.35 / 0.06 %; max single-neuron effect ≈ 4 %",
    "A few FAFB deletions sever whole sensory sources — peripheral single-gateway cells, FAFB-only (0/45 in MCNS)",
    "Whole cell types: nearly all superadditive but barely (median 1.02–1.14×); 14/75 beat matched random sets at z>2 (`VA2_adPN` 3.9×, z=+20)",
    "Cross-modal six: `PVLP076`, `AVLP080` ×2, `CB0677`, `PS124`, `CB0676`. None in all three senses",
]); footer(prs.slides[-1], "B3", "Backup B3")

backup("Metric dependence and ground truth", [
    "Betweenness vs synapse-count top-50 overlap: olf 23/50, legacy-visual 26/50, **mechanosensory 6/50** (chance ≈ 0.3) — mech is the pathway that needed the metric",
    "Flow metrics (Bates 2025 influence; probabilistic traversal) agree with each other (ρ 0.77–0.89), not with betweenness (ρ ≈ 0–0.3)",
    "Current-flow betweenness infeasible (>24 h on smallest graph) — reported as a negative feasibility result",
    "Literature screen: all *necessity* evidence (T4, Mi9 → motion-blind) is in the **control** arm; `lLN2T_c` silencing showed no odour-coding effect — consistent with the deletion result",
]); footer(prs.slides[-1], "B4", "Backup B4")

backup("Replication detail (FAFB vs MaleCNS, 150 sampled sources, cell-type max)", [
    "Olfactory 9/25 (exp 0.52, p 2e-10), ρ 0.67: `lLN2T_c`, `v2LN30`, `il3LN6`, `MZ_lv2PN`, `AL-AST1`, `DM1_lPN` …",
    "Mechanosensory 12/25 (exp 0.28), ρ 0.71: `DNg62`, `DNge132`, `DNge027`, `DNg35`, `PS124`, `PVLP076` …",
    "Visual 13/25 (exp 0.27), ρ 0.70: `Am1`, `H2`, `LPT26`, `LT62`, `LT79`, `PVLP011`, `PS124` …",
    "Degree baseline: 20/18/19 of 25, ρ ≈ 0.9 — the *pathway* is conserved, not extra betweenness information",
    "Only mechanosensory reaches VNC motor neurons within 2+2 hops in MCNS — justifies the descending endpoint",
], fig="talk_fig_replication.png"); footer(prs.slides[-1], "B5", "Backup B5")

out = os.path.join(ROOT, T["OUT"]); prs.save(out); print("saved", out)
